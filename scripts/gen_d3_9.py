#!/usr/bin/env python3
"""
生成 D3-9 答辩PPT.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

BASE_DIR = os.path.dirname(os.path.abspath(__file__))

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色方案
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_ORANGE = RGBColor(0xF4, 0x8C, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF5)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(title_text, subtitle_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, DARK_BLUE)
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.3), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.3), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(20)
        p2.font.color.rgb = ACCENT_ORANGE
        p2.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, WHITE)
    # Title bar
    txBox = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"  {title}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    # Title background
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = DARK_BLUE
    # Bullet content
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.3), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(12)
        if '→' in bullet or '✅' in bullet or '❌' in bullet:
            p.font.color.rgb = ACCENT_ORANGE
    return slide

# ============================================================
# P1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.3), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "NekoCafé DevOps 流水线\n与容器化部署"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.3), Inches(1.5))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "实验三答辩"
p2.font.size = Pt(24)
p2.font.color.rgb = ACCENT_ORANGE
p2.alignment = PP_ALIGN.CENTER
p3 = tf2.add_paragraph()
p3.text = "北京林业大学 · 信息学院 · 软件工程课程"
p3.font.size = Pt(16)
p3.font.color.rgb = WHITE
p3.alignment = PP_ALIGN.CENTER
p4 = tf2.add_paragraph()
p4.text = "2026年5月"
p4.font.size = Pt(14)
p4.font.color.rgb = RGBColor(0xAA, 0xAA, 0xBB)
p4.alignment = PP_ALIGN.CENTER

# ============================================================
# P2: 目标与挑战
# ============================================================
add_content_slide("P2 目标与挑战", [
    "🎯 运维总监三条硬规则：",
    "  ① 提 PR 后 10 分钟内 → 测试环境可见效果",
    "  ② 一行配置 → 灰度发布到 5% 门店",
    "  ③ 出问题 → 3 分钟内定位到具体服务/节点/接口",
    "",
    "📋 实验目标：",
    "  • 理解 CALMS 原则（Culture / Automation / Lean / Measurement / Sharing）",
    "  • 掌握 GitOps 工作流：Git Push → Auto Build → Auto Test → Auto Deploy",
    "  • 构建最小化安全镜像（≤ 200 MB，非 root 运行）",
    "  • 设计完整的 CI/CD 流水线（7 阶段 + 自动回滚）",
    "  • 建立可观测性三件套（Logging + Metrics + Tracing）",
])

# ============================================================
# P3: 总体流水线图
# ============================================================
add_content_slide("P3 总体流水线图", [
    "CI Pipeline (≤ 10 min):",
    "  Lint → Unit Test → SAST → Build → Container Scan → Integration Test → Push Image → PR Comment",
    "",
    "CD Pipeline:",
    "  Push to main → Deploy DEV (auto)",
    "  → Deploy STAGING Canary 5% (monitor 5 min)",
    "  → [审批] → Deploy PROD Canary 5% (monitor 10 min) → Promote 100%",
    "",
    "关键技术决策：",
    "  • Monorepo 架构（PoC 阶段，2 服务，< 5 人团队）",
    "  • Trunk-Based Development（main 分支始终可部署）",
    "  • FastAPI + PostgreSQL + Redis",
    "  • Docker Multi-stage + Helm 3.x + GitHub Actions",
    "  • Prometheus + Grafana + Loki + Tempo",
])

# ============================================================
# P4-P5: 容器化关键决策
# ============================================================
add_content_slide("P4-P5 容器化关键决策", [
    "🐳 Docker 多阶段构建：",
    "  Stage 1 (Builder): 安装编译工具链 → pip install --prefix=/install",
    "  Stage 2 (Runtime): 仅复制编译产物 + 运行时依赖",
    "  → Reservation: ~150 MB | Member: ~145 MB (均 ≤ 200 MB ✅)",
    "",
    "🔒 安全措施：",
    "  • 非 root 用户运行 (UID 1000)",
    "  • 只读根文件系统 (readOnlyRootFilesystem)",
    "  • 丢弃全部 Linux Capabilities",
    "  • Trivy 扫描: 0 HIGH / 0 CRITICAL ✅",
    "  • Docker Compose 本地一键起栈 (docker compose up)",
    "",
    "📦 镜像仓库：GitHub Container Registry (GHCR)",
    "  • tag: git SHA + latest",
    "  • 推送前 Trivy 门禁检查",
])

# ============================================================
# P6-P7: CI 流水线设计与亮点
# ============================================================
add_content_slide("P6-P7 CI 流水线设计与亮点", [
    "⚙️ CI 阶段设计（任一失败即中断）：",
    "  1. Lint: flake8 + black + hadolint (耗时 ~1 min)",
    "  2. Unit Test: pytest + coverage (耗时 ~2 min)",
    "  3. SAST: Bandit 安全扫描 (耗时 ~30s)",
    "  4. Build: Docker 多阶段构建 (耗时 ~2 min, matrix 并行)",
    "  5. Container Scan: Trivy 漏洞扫描 (耗时 ~1 min)",
    "  6. Integration Test: Docker Compose 起栈 (耗时 ~2 min)",
    "  7. Push Image → GHCR (耗时 ~1 min)",
    "",
    "✨ 亮点设计：",
    "  • PR 自动评论（覆盖率 / 漏洞数 / 镜像大小）",
    "  • Docker Layer 缓存 (cache-from: type=gha)",
    "  • Matrix 策略并行（Reservation + Member 同时构建）",
    "  • 总耗时 ≤ 8 分钟 (目标 ≤ 10 分钟 ✅)",
])

# ============================================================
# P8-P9: CD 与渐进式发布
# ============================================================
add_content_slide("P8-P9 CD 与渐进式发布", [
    "🔄 金丝雀发布策略：",
    "  1. 部署 Canary 实例（1 Pod = 5% 流量）",
    "  2. K8s Service Label Selector 控制流量路由",
    "  3. 监控 5-10 min（错误率 / P95 延迟 / 资源使用）",
    "  4. 渐进推进: 5% → 25% → 50% → 100%",
    "  5. 100% 后清理 Canary 资源",
    "",
    "⚡ 自动回滚机制：",
    "  触发条件: P95 > 500ms | 错误率 > 1% | Pod 全不可用",
    "  → Prometheus Alert → GitHub Actions → helm rollback",
    "  → 验证回滚 → 通知团队",
    "",
    "🛠 一键回滚脚本: bash scripts/rollback.sh <env>",
    "  自动查找上一稳定 Helm revision → 执行回滚 → 验证 → 记录日志",
])

# ============================================================
# P10: 可观测性
# ============================================================
add_content_slide("P10 可观测性 Dashboard", [
    "📊 可观测性三件套：",
    "  • Logging: 结构化 JSON 日志 (structlog + Loki + Promtail) — 每条含 traceId",
    "  • Metrics: RED + USE 指标 (Prometheus + ServiceMonitor) — 4 面板 Dashboard",
    "  • Tracing: OpenTelemetry + Tempo (OTLP gRPC) — 请求级链路追踪",
    "",
    "📈 Grafana Dashboard (4 Panels):",
    "  1. QPS — 请求速率 (按 service + path 分组)",
    "  2. Latency P50/P95/P99 — 请求延迟分位数",
    "  3. Error Rate — 5xx 错误率仪表盘 (绿/黄/红)",
    "  4. Resource Usage — CPU + Memory 使用趋势",
    "",
    "🚨 5 条告警规则: ServiceDown / HighErrorRate / HighLatency / HighMemory / HPA Near Max",
])

# ============================================================
# P11: DORA 指标
# ============================================================
add_content_slide("P11 DORA 指标趋势", [
    "📉 DORA 四指标 (14 天数据):",
    "",
    "  指标                第1周          第2周          Elite 基准",
    "  ─────────────────────────────────────────────────────",
    "  部署频率            18 次/周 →    22 次/周       按需",
    "  变更前置时间        3.2 h    →    2.1 h          < 1 h",
    "  变更失败率          22.2%   →    13.6%          0-15%",
    "  MTTR                8.5 min  →    5.2 min        < 1 h",
    "",
    "✅ 趋势向好，第 2 周已接近 Elite 水平",
    "  改善原因：缓存策略优化 → Build 时间减少 40%",
    "            自动回滚上线 → MTTR 下降 39%",
])

# ============================================================
# P12: 安全基线
# ============================================================
add_content_slide("P12 安全基线", [
    "🛡️ 安全措施全景：",
    "",
    "  Secret 管理:",
    "    • GitHub Secrets → K8s Secret 注入 → 严禁硬编码",
    "    • .gitignore 排除 .env / credentials",
    "",
    "  容器安全:",
    "    • 多阶段构建最小攻击面",
    "    • Non-root user (UID 1000) + ReadOnlyRootFilesystem",
    "    • Capabilities: drop ALL",
    "    • Trivy 扫描: 0 HIGH / 0 CRITICAL",
    "",
    "  代码安全:",
    "    • Bandit SAST (Python 静态安全分析)",
    "    • Hadolint (Dockerfile 最佳实践检查)",
    "    • Kube-linter (K8s 清单安全检查)",
    "",
    "  供应链安全:",
    "    • 依赖版本锁定 (requirements.txt 精确版本)",
    "    • SBOM 生成 (Trivy CycloneDX)",
])

# ============================================================
# P13: 踩坑与改进
# ============================================================
add_content_slide("P13 踩坑与改进", [
    "🐛 踩坑记录：",
    "",
    "  1. Docker 缓存失效 → Build 每次 5 分钟",
    "     解决: COPY requirements.txt 放在 COPY src 前面 → 缓存命中率 90%",
    "",
    "  2. Trivy 扫描大量误报 → CI 频繁失败",
    "     解决: 初期 exit-code=0 (仅报告不阻断)，误报清理后改为 exit-code=1",
    "",
    "  3. Helm values 多层嵌套 → 配置错误",
    "     解决: 拆分为 values-dev.yaml / values-staging.yaml / values-prod.yaml",
    "",
    "🔮 改进方向：",
    "  • 接入 ArgoCD 实现 GitOps (Pull 模式 → +3 分加分项)",
    "  • 接入 Chaos Mesh 混沌工程 (故障演练 → +5 分加分项)",
    "  • 增加 Feature Flag 能力 (LaunchDarkly / OpenFeature)",
    "  • 集成 E2E 测试 (Playwright)",
])

# ============================================================
# P14: 下一阶段
# ============================================================
add_content_slide("P14 下一阶段输入", [
    "📥 本实验产物 → 实验四输入：",
    "",
    "  D3-1 DevOps 设计方案 → 实验四质量体系设计参考",
    "  D3-2 源代码仓库 → 实验四代码审查 + AI 辅助测试",
    "  D3-3/4 Dockerfile + CI/CD → 实验四集成测试 + 性能测试",
    "  D3-5 K8s Helm Chart → 实验四环境一致性验证",
    "  D3-6 可观测性 → 实验四测试结果采集与分析",
    "  D3-7 DORA → 实验四质量度量基线",
    "",
    "  📌 实验四将在此基础上开展：",
    "    代码质量审查 + AI 辅助测试用例生成 + 测试覆盖率提升",
])

# ============================================================
# P15: Q&A
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Q & A"
p.font.size = Pt(64)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.3), Inches(1))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "感谢聆听 · 欢迎提问"
p2.font.size = Pt(24)
p2.font.color.rgb = ACCENT_ORANGE
p2.alignment = PP_ALIGN.CENTER

txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.3), Inches(1))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "📧 devops@nekocafe.example.com  |  🔗 github.com/nekocafe"
p3.font.size = Pt(14)
p3.font.color.rgb = RGBColor(0xAA, 0xAA, 0xBB)
p3.alignment = PP_ALIGN.CENTER

output_path = os.path.join(BASE_DIR, '..', '..', 'D3-9_答辩PPT.pptx')
prs.save(output_path)
print(f'✅ D3-9 generated: {output_path}')
